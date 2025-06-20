import re, os, praw, requests, pytz, time, json, uuid
from django.shortcuts import render, get_object_or_404
from django.http import HttpResponse
from django.utils.timezone import datetime
from django.utils import timezone
from .chat_completion import generate_prediction, generate_recap, generate_tweet, generate_parlay, generate_news, generate_videoText, generate_prop, generate_odds, generate_slide_content, get_results
from .image_generation import generate_image, createImagePrompt, generate_image2
from praw.models import InlineImage
from dotenv import load_dotenv
from datetime import datetime as dtdt
from django.http import JsonResponse
import facebook as fb
import tweepy
import openai
import psycopg2
import base64, unicodedata
from .models import Predictions, Recaps, Parlays, Props
from pptx import Presentation
from pptx.util import Inches
from index_now import submit_url_to_index_now, IndexNowAuthentication


load_dotenv(override=True)


ODDSAPI_API_KEY=os.environ.get("ODDSAPI_API_KEY")
consumer_key=os.environ.get('TWITTER_API_KEY')
consumer_secret=os.environ.get('TWITTER_API_SECRET')
bearer_token=os.environ.get('TWITTER_BEARER_TOKEN')
access_token_secret=os.environ.get('TWITTER_ACCESS_TOKEN_SECRET')
access_token=os.environ.get('TWITTER_ACCESS_TOKEN')
FACEBOOK_ACCESS_TOKEN=os.environ.get('FACEBOOK_ACCESS_TOKEN')
DB_NAME=os.environ.get('DB_NAME')
DB_USER=os.environ.get('DB_USER')
DB_PASSWORD=os.environ.get('DB_PASSWORD')
DB_HOST=os.environ.get('DB_HOST')

INauthentication = IndexNowAuthentication(
    host="https://www.gptsportswriter.com",
    api_key="c41a1070cbfa4d5ea66773fc0519716c",
    api_key_location="https://www.gptsportswriter.com/c41a1070cbfa4d5ea66773fc0519716c.txt",
)


ept = pytz.timezone('US/Eastern')
utc = pytz.utc
# str format
fmt = '%Y-%m-%d %H:%M:%S %Z%z'

reddit = praw.Reddit(
    client_id=os.environ.get("REDDIT_CLIENT_ID"),
    client_secret=os.environ.get("REDDIT_CLIENT_SECRET"),
    user_agent="GPTSportsWriter by u/GPTSportsWriter",
    username=os.environ.get("REDDIT_USERNAME"),
    password=os.environ.get("REDDIT_PASSWORD")
)

subreddit = reddit.subreddit("gptsportswriter")

tweepy_auth = tweepy.OAuth1UserHandler(
    "{}".format(os.environ.get("TWITTER_API_KEY")),
    "{}".format(os.environ.get("TWITTER_API_SECRET")),
    "{}".format(os.environ.get("TWITTER_ACCESS_TOKEN")),
    "{}".format(os.environ.get("TWITTER_ACCESS_TOKEN_SECRET")),
)

def recent_predictions(request):
    now = timezone.now()
    twenty_fours_hours_ago = now - timezone.timedelta(hours=48)
    #data = Predictions.objects.filter(created_at__gte=twenty_fours_hours_ago)
    data = Predictions.objects.filter(created_at__gte=twenty_fours_hours_ago).order_by('-created_at').values('id', 'title', 'created_at', 'slug', 'sport_key', 'tweet_text')
    for item in data:
        item['title'] = item['title'].replace("Prediction: ", "")
        item['tweet_text'] = item['tweet_text'].partition(":")[2]
    return render(request, 'predictions/recent_predictions.html', {'data': data})

'''
def prediction_results(request):
    now = timezone.now()
    #twenty_fours_hours_ago = now - timezone.timedelta(hours=36)
    #yesterday = str(twenty_fours_hours_ago)[:10]
    yesterday = "2025-06-11"
    #print(yesterday)
    #data = Predictions.objects.filter(created_at__gte=twenty_fours_hours_ago)
    data = Predictions.objects.filter(title__contains=yesterday).order_by('-created_at').values('id', 'title', 'created_at', 'slug', 'sport_key', 'tweet_text', 'content', 'results')
    for item in data:
        item['title'] = item['title'].replace("Prediction: ", "")
        if item['results'] is None:
         item['results'] = get_results(item['content'], item['title'], item['id'], item['sport_key'])
    totalBets = Predictions.objects.filter(won__isnull=False).count()
    won = Predictions.objects.filter(title__contains=yesterday, won=True).count()
    lost = Predictions.objects.filter(title__contains=yesterday, won=False).count()
    unknown = Predictions.objects.filter(title__contains=yesterday, won=None).count()
    return render(request, 'predictions/prediction_results.html', {'data': data})
'''

def prediction_results(request):
    now = timezone.now()
    thirty_hours_ago = now - timezone.timedelta(hours=30)
    yesterday = str(thirty_hours_ago)[:10]
    data = Predictions.objects.filter(title__contains=yesterday).order_by('-created_at').values('id', 'title', 'created_at', 'slug', 'sport_key', 'tweet_text', 'content', 'results', 'won')
    
    summary = {
        'total_bets': Predictions.objects.filter(title__contains=yesterday).count(),
        'won': Predictions.objects.filter(title__contains=yesterday, won=True).count(),
        'lost': Predictions.objects.filter(title__contains=yesterday, won=False).count(),
        'unknown': Predictions.objects.filter(title__contains=yesterday, won=None).count(),
        'win_rate': 0,
        'loss_rate': 0
    }
    
    if summary['total_bets'] > 0:
        summary['win_rate'] = (summary['won'] / summary['total_bets']) * 100
        summary['loss_rate'] = (summary['lost'] / summary['total_bets']) * 100
    
    for item in data:
        item['title'] = item['title'].replace("Prediction: ", "")
        if item['results'] is None:
            item['results'] = get_results(item['content'], item['title'], item['id'], item['sport_key'])
    
    return render(request, 'predictions/prediction_results.html', {'data': data, 'summary': summary})

def recent_parlays(request):
    now = timezone.now()
    twenty_fours_hours_ago = now - timezone.timedelta(hours=48)
    #data = Parlays.objects.filter(created_at__gte=twenty_fours_hours_ago)
    data = Parlays.objects.filter(created_at__gte=twenty_fours_hours_ago).order_by('-created_at').values('title', 'content', 'created_at', 'slug', 'sport_key', 'tweet_text')
    for item in data:
        item['title'] = item['title'].replace("Parlay: ", "")
        item['tweet_text'] = item['tweet_text'].partition(":")[2]
    return render(request, 'predictions/recent_parlays.html', {'data': data})

def recent_props(request):
    now = timezone.now()
    twenty_fours_hours_ago = now - timezone.timedelta(hours=48)
    #data = Props.objects.filter(created_at__gte=twenty_fours_hours_ago)
    data = Props.objects.filter(created_at__gte=twenty_fours_hours_ago).order_by('-created_at').values('title', 'content', 'created_at', 'slug', 'sport_key', 'tweet_text')
    for item in data:
        item['title'] = item['title'].replace("Prop Bets: ", "")
        item['tweet_text'] = item['tweet_text'].partition(":")[2]
    return render(request, 'predictions/recent_props.html', {'data': data})

def recent_recaps(request):
    now = timezone.now()
    twenty_fours_hours_ago = now - timezone.timedelta(hours=48)
    #data = Recaps.objects.filter(created_at__gte=twenty_fours_hours_ago)
    data = Recaps.objects.filter(created_at__gte=twenty_fours_hours_ago).order_by('-created_at').values('title', 'content', 'created_at', 'slug', 'sport_key', 'tweet_text')
    for item in data:
        item['title'] = item['title'].replace("Recap: ", "")
        item['tweet_text'] = item['tweet_text'].partition(":")[2]
    return render(request, 'predictions/recent_recaps.html', {'data': data})

def prediction_detail(request, slug):
    article = get_object_or_404(Predictions, slug=slug)
    res = ""
    #latest_odds = generate_odds(article.sport_key + " " + article.title, res, article.id, article.sport_key)
    #print(latest_odds)
    context = {
        'article': article,
        #'update': latest_odds.replace("\n", "<br/>"),
    }
    return render(request, 'predictions/prediction_detail.html', context)

def article_detail(request, slug):
    article = get_object_or_404(Predictions, slug=slug)
    res = ""
    #latest_odds = generate_odds(article.sport_key + " " + article.title, res, article.id, article.sport_key)
    #print(latest_odds)
    context = {
        'article': article,
        #'update': latest_odds.replace("\n", "<br/>"),
    }
    return render(request, 'predictions/article_detail.html', context)

def prop_detail(request, slug):
    article = get_object_or_404(Props, slug=slug)
    return render(request, 'predictions/prop_detail.html', {'article': article})

def recap_detail(request, slug):
    article = get_object_or_404(Recaps, slug=slug)
    return render(request, 'predictions/recap_detail.html', {'article': article})

def parlay_detail(request, slug):
    article = get_object_or_404(Parlays, slug=slug)
    return render(request, 'predictions/parlay_detail.html', {'article': article})

# send Tweet
def sendTweet(text, match, file, link):
    tweepy_api = tweepy.API(tweepy_auth)
    #rate_limit_status = tweepy_api.rate_limit_status()
    #rate_limit_status = json.dumps(rate_limit_status)
    #print(rate_limit_status)
    #tweet_endpoint_limit = rate_limit_status['resources']['/tweets&post']['remaining']
    #print(tweet_endpoint_limit)
    #if tweet_endpoint_limit > 0:
    client = tweepy.Client(
    consumer_key=consumer_key,
    consumer_secret=consumer_secret,
    access_token=access_token,
    access_token_secret=access_token_secret
        #wait_on_rate_limit=True
    )

    
    tweetText = createTweet(text)
    tweetText = match + ": " + tweetText
    tweetText = tweetText + " Want more? Visit " + link
    #print(tweetText)

    post = tweepy_api.simple_upload(file)
    text = str(post)
    media_id = re.search("media_id=(.+?),", text).group(1)
    

    # Post Tweet
    #rate_limit_status = client.rate_limit_status()
    #print(rate_limit_status)
    #tweet_endpoint_limit = rate_limit_status['resources']['statuses']['/statuses/update']['remaining']
    #print(tweet_endpoint_limit)
    #if tweet_endpoint_limit > 0:
    try:
        response = client.create_tweet(text=tweetText, media_ids=[media_id])
    except:
        print(tweetText)
    #print(response.data['id'])
    return tweetText
    #else:
        
        #print("Rate limit exceeded. Please wait before tweeting again.")
    #print("twitter: " + response)

def openAITTS(text):
    speech_file_path = "speech.mp3"
    response = openai.audio.speech.create(
        model="tts-1",
        voice="echo",
        speed=1,
        input=text
        )
    response.stream_to_file(speech_file_path)



#def generate_image(prompt):
#    response = openai.images.generate(prompt=prompt[:1000], n=1, size="512x512")
#    image_url = response.data[0].url
#    print(image_url)
#    image_data = requests.get(image_url).content
#    return BytesIO(image_data)



def parse_and_create_ppt(text, topic):
    prs = Presentation()
    slides_data = text.strip().split("\n\n")

    for i, slide in enumerate(slides_data):
        lines = slide.strip().split("\n")
        title = lines[0].split(": ", 1)[1]
        bullets = [line.split(": ", 1)[1] if ": " in line else line for line in lines[1:]]
        content_text = '\n'.join(bullets)

        slide_layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)
        slide_obj.shapes.title.text = title
        slide_obj.placeholders[1].text = content_text

        try:
            img = generate_image2(f"{title}, {bullets[0]}")
            left = Inches(5.5)
            top = Inches(1.5)
            height = Inches(3.5)
            slide_obj.shapes.add_picture(img, left, top, height=height)
        except Exception as e:
            print(f"Image skipped: {e}")

    output_path = f"{topic[:40].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_path)
    return output_path



def createTweet(text):
    tweetText = generate_tweet(text)
    tweetText = tweetText.replace('"', '')
    tweetText = tweetText
    #print(tweetText)
    return tweetText

def fbPost(text, title, file):
    #gptsportswriterapi=fb.GraphAPI("EAAK2IMS7Bb8BO2pYq60uSiYwuJwvMf8ZCcr02yuGwZCFoZCVtzhattkaq9QRQoGAGpZBvaISiZBxutz7hZAvfgDD9VD7T5AEUicF1VztjGtR1WzTFR4cr60YHE7zQZBWbjWwapUASOctQTXY9PPnUeqlJjEhBTdADDYjS22VlM4eZCKiGULGB6lHcDZBZBz8EF")
    #gptsportswriterapi.put_object("me","feed",message=text,link="https://www.gptsportswriter.com")
    #
    postBody = "Prediction" + title + "\n" + "by https://www.gptsportswriter.com" + "\n" + text + "\n\n" + "BetUS - Get 125% Bonus On Your First 3 Deposits, click on the link below..." + "\n" + "https://record.revmasters.com/_8ejz3pKmFDuMKNOJN2Xw7mNd7ZgqdRLk/1/"
    #print(FACEBOOK_ACCESS_TOKEN)
    gptsportswriterapi=fb.GraphAPI(FACEBOOK_ACCESS_TOKEN)
    gptsportswriterapi.put_photo(open(file,'rb'), message=postBody)
    #print(response_photo)
    #photoJson = json.loads(response_photo)
    #photo_id = photoJson[0]['id']
    #print(photo_id)
    #text = "by https://www.gptsportswriter.com " + text
    #gptsportswriterapi.put_object(parent_object="me",connection_name="feed",message=text,link="https://www.gptsportswriter.com",photo_id=photo_id)
    #gptsportswriterapi.put_object("me","feed",message=text,link="https://www.gptsportswriter.com")
    #facebook_access_token_1=(os.environ.get('FACEBOOK_ACCESS_TOKEN'))
    #page_id_1=(os.environ.get('FACEBOOK_PAGE_ID'))
    #msg = title + "\n" + text
    #post_url = 'https://graph.facebook.com/{}/feed`.format(page_id_1)'

def getSports():
    sports = []
    excluded_leagues = ['americanfootball_nfl_preseason','americanfootball_nfl','americanfootball_ncaaf']
    sport = requests.get(f"https://api.the-odds-api.com/v4/sports/?apiKey={ODDSAPI_API_KEY}")
    sport = sport.json()
    for i in range(len(sport)):
        if sport[i]['has_outrights'] == False and sport[i]['key'] not in excluded_leagues:
            sports.append(sport[i]['key'])
            #print(sport[i]['key'])
    return(sports)

def getLeagues():
    excluded_leagues = ['americanfootball_nfl_preseason','americanfootball_nfl','americanfootball_ncaaf']
    leagues = []
    sport = requests.get(f"https://api.the-odds-api.com/v4/sports/?apiKey={ODDSAPI_API_KEY}")
    sport = sport.json()
    for i in range(len(sport)):
        print(sport[i]['key'])
        if sport[i]['has_outrights'] == False and sport[i]['key'] not in excluded_leagues:
            leagues.append(sport[i]['description'])
            #print(sport[i]['key'])
    leagues = [i for n, i in enumerate(leagues) if i not in leagues[:n]]
    #leagues.append("March Madness 2025")
    #leagues.append("Sportsbetting Money Management")
    return(leagues)
            
def ajax_handler(request,sport):
    games = []
    dataMatch = requests.get(f"https://api.the-odds-api.com/v4/sports/{sport}/odds/?apiKey={ODDSAPI_API_KEY}&regions=us&markets=h2h&bookmakers=draftkings,fanduel,betus")
    dataMatch = dataMatch.json()
    #print("predictions: " + str(dataMatch))
    for i in range(len(dataMatch)):
        try:
            t = dataMatch[i]['commence_time']
        except:
            t = "2024-02-25 12:00:00-05:00"
        utcTime = dtdt(int(t[0:4]), int(t[5:7]), int(t[8:10]), int(t[11:13]), int(t[14:16]), int(t[17:19]), tzinfo=utc)
        esTime = utcTime.astimezone(ept)
        now = timezone.now()
        #one_week_from_now = now + timezone.timedelta(hours=168)
        #print(str(t) + ":" + str(one_week_from_now))
        dataMatch[i]['id'] + "test"
        games.append(dataMatch[i]['away_team'] + " VS " + dataMatch[i]['home_team'] + " " + str(esTime) + ":" + dataMatch[i]['id'])
        
    return JsonResponse({'games': games})

def ajax_handlerb(request,sport):
    games = []
    dataMatch = requests.get(f"https://api.the-odds-api.com/v4/sports/{sport}/scores/?apiKey={ODDSAPI_API_KEY}&daysFrom=2")
    dataMatch = dataMatch.json()
    #print("recaps: " + str(dataMatch))
    for i in range(len(dataMatch)):
        try:
            t = dataMatch[i]['commence_time']
        except:
            t = "2024-02-25 12:00:00-05:00"
        utcTime = dtdt(int(t[0:4]), int(t[5:7]), int(t[8:10]), int(t[11:13]), int(t[14:16]), int(t[17:19]), tzinfo=utc)
        esTime = utcTime.astimezone(ept)
        if dataMatch[i]['completed'] == True:
            games.append(dataMatch[i]['away_team'] + " VS " + dataMatch[i]['home_team'] + " " + str(esTime) + ":" + dataMatch[i]['id'])
        
    return JsonResponse({'games': games})

def create_link(page, title):
    baseURL = f"https://www.gptsportswriter.com/{page}-detail/"
    normalized = unicodedata.normalize('NFD', title)
    nlink = u"".join([c for c in normalized if not unicodedata.combining(c)])
    nlink = nlink.lower()
    nlink = nlink.replace(":","")
    nlink = nlink.replace(" ", "-")
    nlink = nlink.replace("(", "")
    nlink = nlink.replace(")", "")
    nlink = nlink.replace("&", "")
    nlink = nlink.replace(".", "")
    nlink = nlink.replace("'", "")
    nlink = nlink[:-1]
    nlink = baseURL + nlink + "/"
    return nlink


# Create your views here.
def home(request):
    return render(request, "predictions/home.html")

def about(request):
    return render(request, "predictions/about.html")

def fbprivacy(request):
    return render(request, "predictions/fbprivacy.html")

def discordTOS(request):
    return render(request, "predictions/discordTOS.html")

def discordPrivacy(request):
    return render(request, "predictions/discordPrivacy.html")

def zoomTOS(request):
    return render(request, "predictions/zoom-tos.html")

def zoomPrivacy(request):
    return render(request, "predictions/zoom-privacy.html")

def zoomSupport(request):
    return render(request, "predictions/zoom-support.html")

def zoomDocumentation(request):
    return render(request, "predictions/zoom-documentation.html")

def sports_betting_money_management(request):
    return render(request, "predictions/sports_betting_money_management.html")

def disclaimer(request):
    return render(request, "predictions/disclaimer.html")

def parlays(request):
    context = {}
    user_input = ""
    sportKey = ""
    sport = ""
    sports = getSports()
    
    if request.method == "GET":
        dataSports = getSports()
        return render(request, "predictions/parlays.html", {'sports': dataSports})
    else:
        if "game" in request.POST:
            user_input += request.POST.get("game") + "\n"
            gameSplit = user_input.split(':')
            gameId=gameSplit[4]
            gameId=gameId.strip()
            match=gameSplit[0]
            sportKey += request.POST.get("sport")
            sport += request.POST.get("sport") + "\n"
            sport = sport.replace('_', " ")
            res = re.split('\s+', match)
            res.remove('VS')
            print(res)

            articles = Parlays.objects.filter(id=gameId)
            if articles:
                for article in articles:
                    imageBytes = get_image_base64(article.gameimg)
                    context = {
                        "user_input": match,
                        "generated_parlay": article.content.replace("\n", "<br/>"),
                        "sports": sports,
                        "image_url":  f"data:;base64,{imageBytes}"
                    }
            else:
                generated_parlay = generate_parlay(sport + " " + match, res, gameId, sportKey)
                image_prompt = createImagePrompt(sport + " " + match)
                #print(image_prompt)
                image_url = generate_image(image_prompt)
                #print(image_url)
                time.sleep(2)
                data = requests.get(image_url).content
                file_name = str(uuid.uuid4()) + ".jpg"
                f = open(file_name, 'wb')
                f.write(data)
                f.close
            
                context = {
                    "user_input": match[:-2],
                    "generated_parlay": generated_parlay.replace("\n", "<br/>"),
                    "image_url": image_url,
                    "sports": sports,
                }

                title = "Parlay: " + match[:-2]
                link = create_link("parlay", title)
                image = InlineImage(path=file_name, caption=title)
                media = {"image1": image}
                selfText = "{image1}" + " by https://www.gptsportswriter.com " + generated_parlay + "\n\nVisit " + link + " for more parlays."
                
                try:
                    subreddit.submit(title, inline_media=media, selftext=selfText) 
                except:
                    print("error submitting reddit post")
        
        
                try:
                    tweetText = sendTweet(generated_parlay, match, file_name, link)
                except:
                    print("error sending tweet")
        
                drawing = open(file_name, 'rb').read()
                parlay = Parlays.objects.create(id=gameId, content=generated_parlay.replace("\n", "<br/>"), gameimg=drawing, title=title, sport_key=sportKey, tweet_text=tweetText)
                try:
                    fbPost(generated_parlay, match, file_name)
                except:
                    print("error posting to FB")
                
                #os.remove(file_name)

        return render(request, "predictions/parlays.html", context)

def topnews(request):
    context = {}
    user_input = ""
    sport = ""
    sports = getLeagues()
    #sports = sports.append("Sports Betting Money Management")
    #sports = ['Baseball MLB','Basketball NCAA','Basketball NBA','Football NCAA','Football NFL','Golf PGA','Ice Hockey NHL','Soccer MLS','Soccer EPL','Tennis','NASCAR Cup Series','Sports Betting Money Management']
    
    if request.method == "GET":
        #dataSports = getLeagues()
        dataSports = sports
        return render(request, "predictions/topnews.html", {'sports': dataSports})
    else:
        if "sport" in request.POST:
            user_input += request.POST.get("sport") + "\n"
            sport += request.POST.get("sport") + "\n"
            sport = sport.replace('_', " ")
            res = re.split('\s+', user_input)
            #print(res)
        
        #print("sport: " + sport)
        generated_news = generate_news(sport, res)
        #print("News: " + generated_news)
        #image_prompt = createImagePrompt(sport)
        #print(image_prompt)
        #image_url = generate_image(image_prompt)
        #print(image_url)
        #time.sleep(2)
        #data = requests.get(image_url).content
        #file_name = str(uuid.uuid4()) + ".jpg"
        #f = open(file_name, 'wb')
        #f.write(data)
        #f.close
            
        context = {
            "user_input": user_input,
            "generated_news": generated_news.replace("\n", "<br/>"),
            "image_url": "",
            "sports": sports,
        }

        title = "Top News: " + user_input
        #image = InlineImage(path=file_name, caption=title)
        #media = {"image1": image}
        selfText = generated_news
        
        try:
            subreddit.submit(title, selftext=selfText)
            #redditURL = subreddit.submit(title, selftext=selfText)
            #redditURL = "https://redd.it/" + str(redditURL)
            #print(redditURL)
        except:
            print("error submitting reddit post")
        
        
        #try:
        #    sendTweet(generated_news, "Top News " + sport )
        #except:
        #    print("error sending tweet")
        

        try:
            fbPost(generated_news, user_input)
        except:
            print("error posting to FB")
        
        #os.remove(file_name)
        return render(request, "predictions/topnews.html", context)

def get_image_base64(image_bytes):
        return base64.b64encode(image_bytes).decode('utf-8')

def predictions(request):
    context = {}
    user_input = ""
    sportKey = ""
    sport = ""
    sports = getSports()

    if request.method == "GET":
        dataSports = getSports()
        return render(request, "predictions/predictions.html", {'sports': dataSports})
    else:
        if "game" in request.POST:
            user_input += request.POST.get("game") + "\n"
            gameSplit = user_input.split(':')
            gameId=gameSplit[4]
            gameId=gameId.strip()
            match=gameSplit[0]
            sportKey += request.POST.get("sport")
            sport += request.POST.get("sport") + "\n"
            sport = sport.replace('_', " ")
            res = re.split('\s+', match)
            res.remove('VS')
            print(res)

            print(gameId)
            articles = Predictions.objects.filter(id=gameId)

            print(articles)
            if articles:
                for article in articles:
                    #print("title: " + article.title)
                    latest_odds = generate_odds(sport + " " + match, res, gameId, sportKey)
                    imageBytes = get_image_base64(article.gameimg)
                    context = {
                        "user_input": match,
                        "generated_prediction": article.content,
                        "sports": sports,
                        "image_url":  f"data:;base64,{imageBytes}",
                        "created_at": article.created_at,
                        "latest_odds": latest_odds.replace("\n", "<br/>")
                    }
            else:
                generated_prediction = generate_prediction(sport + " " + match, res, gameId, sportKey)
                image_prompt = createImagePrompt(sport + " " + match)
                #print(image_prompt)
                image_url = generate_image(image_prompt)
                #print(image_url)
                time.sleep(2)
                try:
                    data = requests.get(image_url).content
                    file_name = str(uuid.uuid4()) + ".jpg"
                    f = open(file_name, 'wb')
                    f.write(data)
                    f.close
                except:
                    print("error creating image")
            
                context = {
                    "user_input": match,
                    "generated_prediction": generated_prediction,
                    "image_url": image_url,
                    "sports": sports,
                    "created_at": "",
                }

                title = "Prediction: " + match[:-2]
                image = InlineImage(path=file_name, caption=title)
                media = {"image1": image}
                link = create_link("prediction", title)
                tweetText = ""
                
                #print(link)
                selfText = "{image1}" + " by https://www.gptsportswriter.com " + generated_prediction + "\n\nVisit " + link + " for more predictions."
                
                #write to database
                #write_to_database(gameId,generated_prediction,"img.jpg",dbTable)
                #drawing = open(file_name, 'rb').read()
                #prediction = Predictions.objects.create(id=gameId, content=generated_prediction.replace("\n", "<br/>"), gameimg=drawing, title=title, sport_key=sportKey)
                
                #content = generate_slide_content(generated_prediction)
                #ppt_path = parse_and_create_ppt(content, generated_prediction)
                #post to reddit
                #post to twitter
                try:
                    tweetText = sendTweet(generated_prediction, match, file_name, link)
                except:
                    print("error sending tweet")

                drawing = open(file_name, 'rb').read()
                print("tweetText:" + tweetText)
                prediction = Predictions.objects.create(id=gameId, content=generated_prediction.replace("\n", "<br/>"), gameimg=drawing, title=title, sport_key=sportKey, tweet_text=tweetText)
                
                articles = Predictions.objects.filter(id=gameId)
                if articles:
                    try:
                        subreddit.submit(title, inline_media=media, selftext=selfText)                                
                    except:
                        print("error submitting reddit post")       
                    #post to facebook
                    try:
                        fbPost(generated_prediction, match, file_name)
                    except:
                        print("error posting to FB")
                else: 
                    pass

                #drawing = open(file_name, 'rb').read()
                #print("tweetText:" + tweetText)
                #prediction = Predictions.objects.create(id=gameId, content=generated_prediction.replace("\n", "<br/>"), gameimg=drawing, title=title, sport_key=sportKey, tweet_text=tweetText)
                #os.remove(file_name)
                try:
                    submit_url_to_index_now(INauthentication, link)
                except:
                    print("faile to submit url")
    
    return render(request, "predictions/predictions.html", context)

def odds(request):
    context = {}
    user_input = ""
    sportKey = ""
    sport = ""
    sports = getSports()

    if request.method == "GET":
        dataSports = getSports()
        return render(request, "predictions/predictions.html", {'sports': dataSports})
    else:
        if "game" in request.POST:
            user_input += request.POST.get("game") + "\n"
            gameSplit = user_input.split(':')
            gameId=gameSplit[4]
            gameId=gameId.strip()
            match=gameSplit[0]
            sportKey += request.POST.get("sport")
            sport += request.POST.get("sport") + "\n"
            sport = sport.replace('_', " ")
            res = re.split('\s+', match)
            res.remove('VS')
            res = res[:len(res)-3]
                   
            
            generated_odds = generate_odds(sport + " " + match, res, gameId, sportKey)
            
            context = {
                "user_input": match,
                "generated_prediction": generated_odds.replace("\n", "<br/>"),
                "image_url": "",
                "sports": sports,
                "created_at": "",
            }

            title = "Odds: " + match[:-2]
            #image = InlineImage(path=file_name, caption=title)
            #media = {"image1": image}
            #selfText = "{image1}" + " by https://www.gptsportswriter.com " + generated_prediction + "\n\nVisit http://www.gptsportswriter.com for more predictions."
                
            #write to database
            #write_to_database(gameId,generated_prediction,"img.jpg",dbTable)
            #drawing = open(file_name, 'rb').read()
            #prediction = Predictions.objects.create(id=gameId, content=generated_prediction.replace("\n", "<br/>"), gameimg=drawing, title=title, sport_key=sportKey)
                
            
    return render(request, "predictions/odds.html", context)

def current_odds(request):
    context = {}
    user_input = ""
    sportKey = ""
    sport = ""
    sports = getSports()

    if request.method == "GET":
        dataSports = getSports()
        return render(request, "predictions/current_odds.html", {'sports': dataSports})
    return render(request, "predictions/current_odds.html", {'sports': ['baseball_mlb','basketball_nba']})

def props(request):
    context = {}
    user_input = ""
    sportKey = ""
    sport = ""
    sports = getSports()
    
    if request.method == "GET":
        dataSports = getSports()
        return render(request, "predictions/props.html", {'sports': dataSports})
    else:
        if "game" in request.POST:
            user_input += request.POST.get("game") + "\n"
            gameSplit = user_input.split(':')
            gameId=gameSplit[4]
            gameId=gameId.strip()
            match=gameSplit[0]
            sportKey += request.POST.get("sport")
            sport += request.POST.get("sport") + "\n"
            sport = sport.replace('_', " ")
            res = re.split('\s+', match)
            res.remove('VS')
            
            articles = Props.objects.filter(id=gameId)
            if articles:
                for article in articles:
                    imageBytes = get_image_base64(article.gameimg)
                    context = {
                        "user_input": match,
                        "generated_prediction": article.content.replace("\n", "<br/>"),
                        "sports": sports,
                        "image_url":  f"data:;base64,{imageBytes}"
                    }
            else:                  
                generated_prop = generate_prop(sport + " " + match, res, gameId, sportKey)
                image_prompt = createImagePrompt(sport + " " + match)
                #print(image_prompt)
                image_url = generate_image(image_prompt)
                #print(image_url)
                time.sleep(2)
                data = requests.get(image_url).content
                file_name = str(uuid.uuid4()) + ".jpg"
                f = open(file_name, 'wb')
                f.write(data)
                f.close
            
                context = {
                    "user_input": match,
                    "generated_prediction": generated_prop.replace("\n", "<br/>"),
                    "image_url": image_url,
                    "sports": sports,
                }

                title = "Prop Bets: " + match[:-2]
                link = create_link("prop", title)
                image = InlineImage(path=file_name, caption=title)
                media = {"image1": image}
                selfText = "{image1}" + " by https://www.gptsportswriter.com " + generated_prop + "\n\nVisit " + link + " for more props."
                
                try:
                    subreddit.submit(title, inline_media=media, selftext=selfText)
                except:
                    print("error submitting reddit post")
        
                try:
                    tweetText = sendTweet(generated_prop, "Prop Bets " + match + " ", file_name, link)
                except:
                    print("error sending tweet")
                
                drawing = open(file_name, 'rb').read()
                prop = Props.objects.create(id=gameId, content=generated_prop.replace("\n", "<br/>") , gameimg=drawing, title=title, sport_key=sportKey, tweet_text=tweetText)
        

                try:
                    fbPost(generated_prop, match, file_name)
                except:
                    print("error posting to FB")
                
                #os.remove(file_name)
        return render(request, "predictions/props.html", context)


def recaps(request):
    context = {}
    user_input = ""
    sportKey = ""
    sports = getSports()
    sport = ""
    
    if request.method == "GET":
        dataSports = getSports()
        return render(request, "predictions/recaps.html", {'sports': dataSports})
    else:
        if "game" in request.POST:
            user_input += request.POST.get("game") + "\n"
            gameSplit = user_input.split(':')
            gameId=gameSplit[4]
            gameId=gameId.strip()
            match=gameSplit[0]
            sportKey += request.POST.get("sport")
            sport += request.POST.get("sport") + "\n"
            sport = sport.replace('_', " ")
            res = re.split('\s+', match)
            res.remove('VS')
            res = res[:len(res)-3]

            articles = Recaps.objects.filter(id=gameId)
            if articles:
                for article in articles:
                    imageBytes = get_image_base64(article.gameimg)
                    context = {
                        "user_input": match,
                        "generated_recap": article.content.replace("\n", "<br/>"),
                        "sports": sports,
                        "image_url":  f"data:;base64,{imageBytes}"
                    }
            else:        
                generated_recap = generate_recap(sport + " " + match, res, gameId, sportKey)
                print(sport)
                image_prompt = createImagePrompt(sport + " " + match)
                #print(image_prompt)
                image_url = generate_image(image_prompt)
                #print(image_url)
                time.sleep(2)
                data = requests.get(image_url).content
                file_name = str(uuid.uuid4()) + ".jpg"
                f = open(file_name, 'wb')
                f.write(data)
                f.close
            
                context = {
                    "user_input": match,
                    "generated_recap": generated_recap.replace("\n", "<br/>"),
                    "image_url": image_url,
                    "sports": sports,
                }

                title = "Recap: " + match[:-2]
                link = create_link("recap", title)
                image = InlineImage(path=file_name, caption=title)
                media = {"image1": image}
                selfText = "{image1}" + " by https://www.gptsportswriter.com " + generated_recap + "\n\nVisit " + link + " for more recaps."
                
                try:
                    subreddit.submit(title, inline_media=media, selftext=selfText)
                except:
                    print("error submitting reddit post")

                try:
                    tweetText = sendTweet(generated_recap, "Recap " + match + " ", file_name, link)
                except:
                    print("error sending tweet")
                
                drawing = open(file_name, 'rb').read()
                recap = Recaps.objects.create(id=gameId, content=generated_recap.replace("\n", "<br/>"), gameimg=drawing, title=title, sport_key=sportKey, tweet_text=tweetText)
    
                try:
                    fbPost(generated_recap, match, file_name)
                except:
                    print("error posting to FB")
                
                #os.remove(file_name)

        return render(request, "predictions/recaps.html", context)
